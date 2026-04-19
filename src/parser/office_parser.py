"""
Word/PPT文档解析器
"""
from pathlib import Path
from typing import Dict, Any
from langchain_core.documents import Document
import docx
from pptx import Presentation


class DocxParser:
    """Word文档解析器"""

    def __init__(self, docx_path: str):
        self.docx_path = Path(docx_path)
        if not self.docx_path.exists():
            raise FileNotFoundError(f"文件不存在: {docx_path}")

    def parse(self) -> Dict[str, Any]:
        """解析Word文档"""
        doc = docx.Document(self.docx_path)

        # 提取文本
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        full_text = "\n\n".join(paragraphs)

        # 提取表格
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        return {
            "file_name": self.docx_path.name,
            "title": self.docx_path.stem,
            "content": full_text,
            "tables": tables,
            "paragraph_count": len(paragraphs)
        }

    def to_document(self, metadata: Dict[str, Any] = None) -> Document:
        """转换为Document"""
        parsed = self.parse()

        base_metadata = metadata or {}
        base_metadata.update({
            "source": str(self.docx_path),
            "file_name": parsed["file_name"],
            "title": parsed["title"],
            "type": "docx"
        })

        # 合并表格内容
        content_parts = [parsed["content"]]
        for i, table in enumerate(parsed["tables"], 1):
            table_text = f"\n\n【表格{i}】\n"
            for row in table:
                table_text += " | ".join(row) + "\n"
            content_parts.append(table_text)

        doc = Document(
            page_content="\n".join(content_parts),
            metadata=base_metadata
        )

        return doc


class PptxParser:
    """PPT文档解析器"""

    def __init__(self, pptx_path: str):
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"文件不存在: {pptx_path}")

    def parse(self) -> Dict[str, Any]:
        """解析PPT文档"""
        prs = Presentation(self.pptx_path)

        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                slides_content.append({
                    "slide_number": i,
                    "content": "\n".join(slide_text)
                })

        full_text = "\n\n---\n\n".join([
            f"【第{s['slide_number']}页】\n{s['content']}"
            for s in slides_content
        ])

        return {
            "file_name": self.pptx_path.name,
            "title": self.pptx_path.stem,
            "content": full_text,
            "slides": slides_content,
            "slide_count": len(slides_content)
        }

    def to_document(self, metadata: Dict[str, Any] = None) -> Document:
        """转换为Document"""
        parsed = self.parse()

        base_metadata = metadata or {}
        base_metadata.update({
            "source": str(self.pptx_path),
            "file_name": parsed["file_name"],
            "title": parsed["title"],
            "type": "pptx",
            "slide_count": parsed["slide_count"]
        })

        doc = Document(
            page_content=parsed["content"],
            metadata=base_metadata
        )

        return doc
